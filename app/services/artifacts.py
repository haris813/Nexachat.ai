from __future__ import annotations

import hashlib
import json
import re
import secrets
import zipfile
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from flask import current_app
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table as OpenpyxlTable
from openpyxl.worksheet.table import TableStyleInfo
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from pypdf import PdfReader
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from ..extensions import db
from ..models import Artifact
from .security import safe_display_name

MIME_TYPES = {
    "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "word": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "powerpoint": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
    "chart": "image/png",
    "audio": "audio/mpeg",
    "conversation": "text/markdown",
}

EXTENSIONS = {
    "excel": ".xlsx",
    "word": ".docx",
    "powerpoint": ".pptx",
    "pdf": ".pdf",
    "chart": ".png",
    "audio": ".mp3",
    "conversation": ".md",
}

THEMES = {
    "modern_light": {"bg": "F7F8FC", "text": "182033", "accent": "6C5CE7", "accent2": "00A8A8"},
    "minimal_dark": {"bg": "11131A", "text": "F2F4FA", "accent": "8B7CFF", "accent2": "48D6CA"},
    "professional_business": {"bg": "F4F7FA", "text": "15263A", "accent": "176B87", "accent2": "64CCC5"},
    "startup_pitch": {"bg": "FFF9F2", "text": "25202A", "accent": "FF5A5F", "accent2": "6C5CE7"},
    "academic": {"bg": "FAFAF7", "text": "202020", "accent": "2F5D50", "accent2": "8C6A3B"},
    "technical_architecture": {"bg": "F3F6FA", "text": "152238", "accent": "315EFB", "accent2": "00A6A6"},
}


def _slug(value: str) -> str:
    result = re.sub(r"[^a-zA-Z0-9]+", "-", value).strip("-").lower()
    return (result or "nexachat-artifact")[:80]


def _safe_cell(value: Any) -> Any:
    if isinstance(value, str) and value.lstrip().startswith(("=", "+", "-", "@")):
        return "'" + value
    return value


def _hex_color(value: str) -> str:
    return value.lstrip("#").upper()


class ArtifactService:
    def __init__(self, owner_id: str, conversation_id: int | None = None, plan_id: str | None = None):
        self.owner_id = owner_id
        self.conversation_id = conversation_id
        self.plan_id = plan_id
        owner_hash = hashlib.sha256(owner_id.encode()).hexdigest()[:16]
        self.output_dir = Path(current_app.config["ARTIFACT_DIR"]) / owner_hash
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def _path(self, title: str, kind: str) -> tuple[Path, str]:
        stored_name = f"{secrets.token_hex(16)}-{_slug(title)}{EXTENSIONS[kind]}"
        return self.output_dir / stored_name, stored_name

    def _record(
        self,
        path: Path,
        stored_name: str,
        title: str,
        kind: str,
        *,
        metadata: dict | None = None,
        preview: dict | None = None,
    ) -> Artifact:
        from .storage import get_artifact_storage

        storage = get_artifact_storage()
        storage_key = storage.save_file(self.owner_id, stored_name, path)
        record = Artifact(
            owner_id=self.owner_id,
            conversation_id=self.conversation_id,
            plan_id=self.plan_id,
            kind=kind,
            display_name=f"{safe_display_name(title)}{EXTENSIONS[kind]}",
            stored_name=stored_name,
            storage_path=storage_key,
            mime_type=MIME_TYPES[kind],
            size_bytes=path.stat().st_size,
            status="ready",
            metadata_json=json.dumps(metadata or {}, ensure_ascii=False, default=str),
            preview_json=json.dumps(preview or {}, ensure_ascii=False, default=str),
        )
        db.session.add(record)
        db.session.commit()
        return record

    def create_excel(
        self,
        title: str,
        rows: list[dict[str, Any]],
        *,
        columns: list[str] | None = None,
        sources: list[dict] | None = None,
        data_date: str | None = None,
    ) -> Artifact:
        if not rows:
            raise ValueError("Excel generation requires at least one data row")
        headers = columns or list(rows[0].keys())
        path, stored_name = self._path(title, "excel")
        workbook = Workbook()
        data_sheet = workbook.active
        data_sheet.title = "Data"
        last_column = get_column_letter(len(headers))
        data_sheet.merge_cells(f"A1:{last_column}1")
        data_sheet["A1"] = title
        data_sheet["A1"].font = Font(name="Aptos Display", size=20, bold=True, color="FFFFFF")
        data_sheet["A1"].fill = PatternFill("solid", fgColor="4F46E5")
        data_sheet["A1"].alignment = Alignment(vertical="center")
        data_sheet.row_dimensions[1].height = 34
        data_sheet.merge_cells(f"A2:{last_column}2")
        data_sheet["A2"] = f"Generated by NexaChat AI on {datetime.now(UTC).date().isoformat()}" + (
            f" | Data date: {data_date}" if data_date else ""
        )
        data_sheet["A2"].font = Font(name="Aptos", size=10, italic=True, color="667085")
        header_row = 4
        for column_index, header in enumerate(headers, start=1):
            cell = data_sheet.cell(row=header_row, column=column_index, value=header)
            cell.font = Font(name="Aptos", bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="252B37")
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        thin = Side(style="thin", color="D0D5DD")
        for row_index, row in enumerate(rows, start=header_row + 1):
            for column_index, header in enumerate(headers, start=1):
                value = _safe_cell(row.get(header))
                cell = data_sheet.cell(row=row_index, column=column_index, value=value)
                cell.font = Font(name="Aptos", size=10, color="1D2939")
                cell.alignment = Alignment(vertical="top", wrap_text=True)
                cell.border = Border(bottom=thin)
                if isinstance(value, (int, float)) and not isinstance(value, bool):
                    lowered = header.lower()
                    if any(
                        term in lowered for term in ("price", "revenue", "cost", "worth", "amount", "sales")
                    ):
                        cell.number_format = "$#,##0.00;[Red]-$#,##0.00"
                    elif any(term in lowered for term in ("percent", "rate", "margin", "%")):
                        cell.number_format = "0.0%"
                    else:
                        cell.number_format = "#,##0.00"
        data_sheet.freeze_panes = f"A{header_row + 1}"
        data_sheet.auto_filter.ref = f"A{header_row}:{last_column}{header_row + len(rows)}"
        table = OpenpyxlTable(displayName="NexaChatData", ref=data_sheet.auto_filter.ref)
        table.tableStyleInfo = TableStyleInfo(
            name="TableStyleMedium2",
            showFirstColumn=False,
            showLastColumn=False,
            showRowStripes=True,
            showColumnStripes=False,
        )
        data_sheet.add_table(table)
        for index, header in enumerate(headers, start=1):
            values = [str(header), *[str(row.get(header, "")) for row in rows[:100]]]
            width = min(max(max(len(value) for value in values) + 2, 12), 46)
            data_sheet.column_dimensions[get_column_letter(index)].width = width

        numeric_indexes = [
            index
            for index, header in enumerate(headers, start=1)
            if any(
                isinstance(row.get(header), (int, float)) and not isinstance(row.get(header), bool)
                for row in rows
            )
        ]
        if numeric_indexes:
            for index in numeric_indexes:
                data_sheet.conditional_formatting.add(
                    f"{get_column_letter(index)}{header_row + 1}:{get_column_letter(index)}{header_row + len(rows)}",
                    ColorScaleRule(
                        start_type="min",
                        start_color="EEF2FF",
                        mid_type="percentile",
                        mid_value=50,
                        mid_color="C7D2FE",
                        end_type="max",
                        end_color="818CF8",
                    ),
                )
            chart_sheet = workbook.create_sheet("Charts")
            chart_sheet.sheet_view.showGridLines = False
            chart_sheet["A1"] = f"{title} - visual summary"
            chart_sheet["A1"].font = Font(name="Aptos Display", size=18, bold=True, color="252B37")
            chart = BarChart()
            chart.type = "bar"
            chart.style = 10
            chart.title = headers[numeric_indexes[0] - 1]
            chart.height = 7.5
            chart.width = 14
            values = Reference(
                data_sheet,
                min_col=numeric_indexes[0],
                min_row=header_row,
                max_row=header_row + min(len(rows), 15),
            )
            categories = Reference(
                data_sheet,
                min_col=1,
                min_row=header_row + 1,
                max_row=header_row + min(len(rows), 15),
            )
            chart.add_data(values, titles_from_data=True)
            chart.set_categories(categories)
            chart.legend = None
            chart_sheet.add_chart(chart, "A3")

        metadata_sheet = workbook.create_sheet("Metadata")
        metadata_rows = [
            ("Artifact", title),
            ("Generated at", datetime.now(UTC).isoformat()),
            ("Data date", data_date or "Not specified"),
            ("Rows", len(rows)),
            ("Generator", "NexaChat AI"),
        ]
        for row in metadata_rows:
            metadata_sheet.append(row)
        metadata_sheet.append([])
        metadata_sheet.append(("Source name", "Source URL", "Retrieved at"))
        for source in sources or []:
            metadata_sheet.append(
                (
                    source.get("title", ""),
                    source.get("url", ""),
                    source.get("retrieved_at", ""),
                )
            )
        for cell in metadata_sheet[7]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="4F46E5")
        metadata_sheet.column_dimensions["A"].width = 34
        metadata_sheet.column_dimensions["B"].width = 78
        metadata_sheet.column_dimensions["C"].width = 28
        metadata_sheet.freeze_panes = "A7"
        workbook.properties.creator = "NexaChat AI"
        workbook.properties.title = title
        workbook.save(path)
        self._validate_excel(path)
        preview = {"columns": headers, "rows": [[row.get(header) for header in headers] for row in rows[:10]]}
        return self._record(
            path,
            stored_name,
            title,
            "excel",
            metadata={"rows": len(rows), "columns": len(headers), "sources": len(sources or [])},
            preview=preview,
        )

    def create_word(
        self,
        title: str,
        sections: list[dict[str, Any]],
        *,
        sources: list[dict] | None = None,
        template: str = "business_report",
    ) -> Artifact:
        path, stored_name = self._path(title, "word")
        document = Document()
        section = document.sections[0]
        section.top_margin = Inches(0.85)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(0.9)
        section.right_margin = Inches(0.9)
        styles = document.styles
        styles["Normal"].font.name = "Aptos"
        styles["Normal"].font.size = Pt(10.5)
        styles["Normal"].paragraph_format.space_after = Pt(7)
        for style_name, size, color in (
            ("Title", 30, RGBColor(37, 43, 55)),
            ("Heading 1", 19, RGBColor(79, 70, 229)),
            ("Heading 2", 14, RGBColor(37, 43, 55)),
        ):
            style = styles[style_name]
            style.font.name = "Aptos Display"
            style.font.size = Pt(size)
            style.font.color.rgb = color
            style.font.bold = True
        if "NexaChat Lead" not in styles:
            lead_style = styles.add_style("NexaChat Lead", WD_STYLE_TYPE.PARAGRAPH)
            lead_style.font.name = "Aptos"
            lead_style.font.size = Pt(12)
            lead_style.font.color.rgb = RGBColor(71, 84, 103)
            lead_style.paragraph_format.space_after = Pt(14)

        document.add_paragraph("NEXACHAT AI  /  INTELLIGENT ARTIFACT", style="Subtitle")
        title_paragraph = document.add_paragraph(title, style="Title")
        title_paragraph.paragraph_format.space_before = Pt(80)
        title_paragraph.paragraph_format.space_after = Pt(18)
        lead = document.add_paragraph(
            f"Professional {template.replace('_', ' ')} generated on {datetime.now(UTC).date().isoformat()}.",
            style="NexaChat Lead",
        )
        lead.paragraph_format.space_after = Pt(24)
        document.add_paragraph(
            "This document was generated from the user request and the evidence listed in References. "
            "Web and uploaded content are treated as untrusted data, never as instructions."
        )
        document.add_page_break()

        document.add_heading("Contents", level=1)
        for index, item in enumerate(sections, start=1):
            document.add_paragraph(
                f"{index}. {item.get('heading') or f'Section {index}'}", style="List Number"
            )
        if sources:
            document.add_paragraph(f"{len(sections) + 1}. References", style="List Number")
        document.add_page_break()

        for item in sections:
            heading = str(item.get("heading") or "Section")
            document.add_heading(heading, level=1)
            body = item.get("body", "")
            if isinstance(body, list):
                for bullet in body:
                    document.add_paragraph(str(bullet), style="List Bullet")
            else:
                for paragraph in str(body).split("\n\n"):
                    if paragraph.strip():
                        document.add_paragraph(paragraph.strip())
            table_rows = item.get("table")
            if table_rows and isinstance(table_rows, list) and isinstance(table_rows[0], dict):
                headers = list(table_rows[0])
                table = document.add_table(rows=1, cols=len(headers))
                table.style = "Light Shading Accent 1"
                table.autofit = False
                for index, header in enumerate(headers):
                    cell = table.rows[0].cells[index]
                    cell.text = header
                    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                    for run in cell.paragraphs[0].runs:
                        run.font.bold = True
                for row in table_rows:
                    cells = table.add_row().cells
                    for index, header in enumerate(headers):
                        cells[index].text = str(row.get(header, ""))
                        cells[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

        if sources:
            document.add_page_break()
            document.add_heading("References", level=1)
            for source in sources:
                paragraph = document.add_paragraph(style="List Number")
                paragraph.add_run(f"{source.get('title', 'Source')}. ").bold = True
                paragraph.add_run(
                    f"{source.get('url', '')} (retrieved {str(source.get('retrieved_at', ''))[:10]})."
                )

        for section in document.sections:
            footer = section.footer.paragraphs[0]
            footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
            footer.add_run("NexaChat AI  •  ")
            self._add_page_number(footer)
        document.core_properties.title = title
        document.core_properties.author = "NexaChat AI"
        document.core_properties.subject = template.replace("_", " ").title()
        document.save(path)
        self._validate_word(path)
        preview = {
            "title": title,
            "sections": [str(item.get("heading") or "Section") for item in sections],
            "source_count": len(sources or []),
        }
        return self._record(
            path,
            stored_name,
            title,
            "word",
            metadata={"template": template, "sections": len(sections), "sources": len(sources or [])},
            preview=preview,
        )

    @staticmethod
    def _add_page_number(paragraph) -> None:
        run = paragraph.add_run()
        begin = OxmlElement("w:fldChar")
        begin.set(qn("w:fldCharType"), "begin")
        instruction = OxmlElement("w:instrText")
        instruction.set(qn("xml:space"), "preserve")
        instruction.text = "PAGE"
        end = OxmlElement("w:fldChar")
        end.set(qn("w:fldCharType"), "end")
        run._r.extend([begin, instruction, end])

    def create_powerpoint(
        self,
        title: str,
        slides: list[dict[str, Any]],
        *,
        sources: list[dict] | None = None,
        theme: str = "modern_light",
        audience: str = "general",
    ) -> Artifact:
        palette = THEMES.get(theme, THEMES["modern_light"])
        path, stored_name = self._path(title, "powerpoint")
        presentation = Presentation()
        presentation.slide_width = PptInches(13.333)
        presentation.slide_height = PptInches(7.5)
        presentation.core_properties.title = title
        presentation.core_properties.author = "NexaChat AI"
        self._ppt_title_slide(presentation, title, f"Prepared for {audience}", palette)
        agenda_titles = [str(item.get("title") or f"Section {index}") for index, item in enumerate(slides, 1)]
        self._ppt_content_slide(
            presentation,
            "Agenda",
            agenda_titles[:8],
            palette,
            section_label="01 / ORIENTATION",
        )
        for index, item in enumerate(slides, start=1):
            bullets = item.get("bullets") or item.get("body") or []
            if isinstance(bullets, str):
                bullets = [part.strip() for part in bullets.split("\n") if part.strip()]
            self._ppt_content_slide(
                presentation,
                str(item.get("title") or f"Insight {index}"),
                [str(value) for value in bullets][:6],
                palette,
                section_label=f"{index + 1:02d} / INSIGHT",
                takeaway=str(item.get("takeaway") or ""),
            )
        if sources:
            source_lines = [
                f"{index}. {source.get('title', 'Source')} — {source.get('url', '')}"
                for index, source in enumerate(sources[:8], start=1)
            ]
            self._ppt_content_slide(
                presentation,
                "Sources",
                source_lines,
                palette,
                section_label=f"{len(slides) + 2:02d} / EVIDENCE",
                body_size=16,
            )
        presentation.save(path)
        self._validate_powerpoint(path)
        preview = {
            "title": title,
            "slides": [
                "Title",
                "Agenda",
                *[str(item.get("title") or "Insight") for item in slides],
                *(["Sources"] if sources else []),
            ],
            "theme": theme,
        }
        return self._record(
            path,
            stored_name,
            title,
            "powerpoint",
            metadata={"slides": len(presentation.slides), "theme": theme, "sources": len(sources or [])},
            preview=preview,
        )

    @staticmethod
    def _ppt_background(slide, palette: dict[str, str]) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = PptRGBColor.from_string(_hex_color(palette["bg"]))

    def _ppt_title_slide(
        self, presentation: Presentation, title: str, subtitle: str, palette: dict[str, str]
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._ppt_background(slide, palette)
        accent = slide.shapes.add_shape(1, PptInches(0.72), PptInches(0.76), PptInches(0.16), PptInches(5.9))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PptRGBColor.from_string(_hex_color(palette["accent"]))
        accent.line.fill.background()
        text_box = slide.shapes.add_textbox(PptInches(1.25), PptInches(1.55), PptInches(10.9), PptInches(3.4))
        frame = text_box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = title[:120]
        paragraph.font.name = "Aptos Display"
        paragraph.font.size = PptPt(50)
        paragraph.font.bold = True
        paragraph.font.color.rgb = PptRGBColor.from_string(_hex_color(palette["text"]))
        paragraph.space_after = PptPt(18)
        sub = frame.add_paragraph()
        sub.text = subtitle
        sub.font.name = "Aptos"
        sub.font.size = PptPt(20)
        sub.font.color.rgb = PptRGBColor.from_string(_hex_color(palette["accent2"]))
        footer = slide.shapes.add_textbox(PptInches(1.25), PptInches(6.65), PptInches(8), PptInches(0.35))
        footer.text_frame.paragraphs[0].text = f"NEXACHAT AI  •  {datetime.now(UTC).date().isoformat()}"
        footer.text_frame.paragraphs[0].font.size = PptPt(11)
        footer.text_frame.paragraphs[0].font.bold = True
        footer.text_frame.paragraphs[0].font.color.rgb = PptRGBColor.from_string(
            _hex_color(palette["accent"])
        )

    def _ppt_content_slide(
        self,
        presentation: Presentation,
        title: str,
        bullets: list[str],
        palette: dict[str, str],
        *,
        section_label: str,
        takeaway: str = "",
        body_size: int = 21,
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._ppt_background(slide, palette)
        label_box = slide.shapes.add_textbox(
            PptInches(0.78), PptInches(0.48), PptInches(3.2), PptInches(0.32)
        )
        label = label_box.text_frame.paragraphs[0]
        label.text = section_label
        label.font.name = "Aptos"
        label.font.bold = True
        label.font.size = PptPt(11)
        label.font.color.rgb = PptRGBColor.from_string(_hex_color(palette["accent"]))
        title_box = slide.shapes.add_textbox(
            PptInches(0.78), PptInches(0.92), PptInches(11.75), PptInches(0.82)
        )
        title_paragraph = title_box.text_frame.paragraphs[0]
        title_paragraph.text = title[:90]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = PptPt(35)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = PptRGBColor.from_string(_hex_color(palette["text"]))
        body_box = slide.shapes.add_textbox(
            PptInches(0.88), PptInches(2.0), PptInches(11.35), PptInches(4.25)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.margin_left = PptInches(0.12)
        body_frame.margin_right = PptInches(0.12)
        body_frame.clear()
        for index, value in enumerate(bullets or ["No supporting points were supplied."]):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = value[:280]
            paragraph.level = 0
            paragraph.font.name = "Aptos"
            paragraph.font.size = PptPt(body_size)
            paragraph.font.color.rgb = PptRGBColor.from_string(_hex_color(palette["text"]))
            paragraph.space_after = PptPt(14)
            paragraph.text = f"•  {paragraph.text}"
        if takeaway:
            takeaway_box = slide.shapes.add_textbox(
                PptInches(0.88), PptInches(6.42), PptInches(11.35), PptInches(0.52)
            )
            takeaway_paragraph = takeaway_box.text_frame.paragraphs[0]
            takeaway_paragraph.text = f"KEY TAKEAWAY  /  {takeaway[:170]}"
            takeaway_paragraph.font.name = "Aptos"
            takeaway_paragraph.font.bold = True
            takeaway_paragraph.font.size = PptPt(13)
            takeaway_paragraph.font.color.rgb = PptRGBColor.from_string(_hex_color(palette["accent2"]))

    def create_pdf(
        self,
        title: str,
        sections: list[dict[str, Any]],
        *,
        sources: list[dict] | None = None,
    ) -> Artifact:
        path, stored_name = self._path(title, "pdf")
        font_name = self._register_pdf_font()
        styles = getSampleStyleSheet()
        styles.add(
            ParagraphStyle(
                "NexaTitle",
                parent=styles["Title"],
                fontName=font_name,
                fontSize=27,
                leading=33,
                textColor=colors.HexColor("#252B37"),
                alignment=TA_CENTER,
                spaceAfter=20,
            )
        )
        styles.add(
            ParagraphStyle(
                "NexaHeading",
                parent=styles["Heading1"],
                fontName=font_name,
                fontSize=16,
                leading=20,
                textColor=colors.HexColor("#4F46E5"),
                spaceBefore=14,
                spaceAfter=8,
            )
        )
        styles["BodyText"].fontName = font_name
        styles["BodyText"].fontSize = 10.5
        styles["BodyText"].leading = 15
        styles["BodyText"].textColor = colors.HexColor("#344054")
        styles["BodyText"].spaceAfter = 8

        def footer(canvas, document):
            canvas.saveState()
            canvas.setFont(font_name, 8)
            canvas.setFillColor(colors.HexColor("#667085"))
            canvas.drawString(0.75 * inch, 0.45 * inch, "NexaChat AI")
            canvas.drawRightString(7.75 * inch, 0.45 * inch, f"Page {document.page}")
            canvas.restoreState()

        report = SimpleDocTemplate(
            str(path),
            pagesize=letter,
            rightMargin=0.75 * inch,
            leftMargin=0.75 * inch,
            topMargin=0.72 * inch,
            bottomMargin=0.72 * inch,
            title=title,
            author="NexaChat AI",
        )
        story: list[Any] = [
            Spacer(1, 1.2 * inch),
            Paragraph(_xml_escape(title), styles["NexaTitle"]),
            Paragraph(
                f"Generated {datetime.now(UTC).date().isoformat()} by NexaChat AI",
                ParagraphStyle(
                    "NexaSubtitle",
                    parent=styles["BodyText"],
                    alignment=TA_CENTER,
                    textColor=colors.HexColor("#667085"),
                ),
            ),
            PageBreak(),
        ]
        for item in sections:
            story.append(Paragraph(_xml_escape(str(item.get("heading") or "Section")), styles["NexaHeading"]))
            body = item.get("body", "")
            values = body if isinstance(body, list) else str(body).split("\n\n")
            for paragraph in values:
                if str(paragraph).strip():
                    story.append(Paragraph(_xml_escape(str(paragraph).strip()), styles["BodyText"]))
            table_rows = item.get("table")
            if table_rows and isinstance(table_rows, list) and isinstance(table_rows[0], dict):
                headers = list(table_rows[0])
                table_data = [[Paragraph(_xml_escape(header), styles["BodyText"]) for header in headers]]
                for row in table_rows[:40]:
                    table_data.append(
                        [
                            Paragraph(_xml_escape(str(row.get(header, ""))), styles["BodyText"])
                            for header in headers
                        ]
                    )
                table = Table(table_data, repeatRows=1, hAlign="LEFT")
                table.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4F46E5")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                            ("FONTNAME", (0, 0), (-1, -1), font_name),
                            ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D0D5DD")),
                            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                            ("LEFTPADDING", (0, 0), (-1, -1), 7),
                            ("RIGHTPADDING", (0, 0), (-1, -1), 7),
                            ("TOPPADDING", (0, 0), (-1, -1), 6),
                            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                        ]
                    )
                )
                story.append(KeepTogether([Spacer(1, 5), table, Spacer(1, 8)]))
        if sources:
            story.append(PageBreak())
            story.append(Paragraph("Sources", styles["NexaHeading"]))
            for index, source in enumerate(sources, start=1):
                story.append(
                    Paragraph(
                        _xml_escape(
                            f"{index}. {source.get('title', 'Source')} - {source.get('url', '')} "
                            f"(retrieved {str(source.get('retrieved_at', ''))[:10]})"
                        ),
                        styles["BodyText"],
                    )
                )
        report.build(story, onFirstPage=footer, onLaterPages=footer)
        self._validate_pdf(path)
        preview = {
            "title": title,
            "sections": [str(item.get("heading") or "Section") for item in sections],
            "source_count": len(sources or []),
        }
        return self._record(
            path,
            stored_name,
            title,
            "pdf",
            metadata={"sections": len(sections), "sources": len(sources or [])},
            preview=preview,
        )

    def create_chart(self, title: str, labels: list[str], values: list[float]) -> Artifact:
        if not labels or len(labels) != len(values):
            raise ValueError("Chart labels and values must be non-empty and have matching lengths")
        path, stored_name = self._path(title, "chart")
        width, height = 1400, 820
        image = Image.new("RGB", (width, height), "#F7F8FC")
        draw = ImageDraw.Draw(image)
        font = ImageFont.load_default()
        draw.text((70, 48), title, fill="#182033", font=font)
        maximum = max(max(values), 1)
        chart_left, chart_top, chart_right, chart_bottom = 120, 135, 1320, 700
        bar_space = (chart_bottom - chart_top) / len(values)
        for index, (label, value) in enumerate(zip(labels, values, strict=True)):
            y = chart_top + index * bar_space + 8
            bar_height = max(16, bar_space - 16)
            bar_width = (chart_right - chart_left) * max(0, value) / maximum
            draw.rounded_rectangle(
                (chart_left, y, chart_left + bar_width, y + bar_height),
                radius=10,
                fill="#6C5CE7",
            )
            draw.text((20, y + 4), str(label)[:18], fill="#344054", font=font)
            draw.text((chart_left + bar_width + 12, y + 4), f"{value:,.2f}", fill="#344054", font=font)
        image.save(path, "PNG")
        return self._record(
            path,
            stored_name,
            title,
            "chart",
            metadata={"points": len(values)},
            preview={"labels": labels, "values": values},
        )

    def create_audio(self, title: str, text: str, synthesizer) -> Artifact:
        path, stored_name = self._path(title, "audio")
        try:
            synthesizer(path)
            if not path.exists() or path.stat().st_size == 0:
                raise ValueError("Speech provider returned an empty audio file")
            return self._record(
                path,
                stored_name,
                title,
                "audio",
                metadata={"characters": len(text)},
                preview={"text": text[:500]},
            )
        except Exception:
            path.unlink(missing_ok=True)
            raise

    def create_conversation_export(self, title: str, markdown: str) -> Artifact:
        path, stored_name = self._path(title, "conversation")
        path.write_text(markdown, encoding="utf-8")
        return self._record(
            path,
            stored_name,
            title,
            "conversation",
            metadata={"characters": len(markdown)},
            preview={"excerpt": markdown[:700]},
        )

    def register_existing(
        self,
        path: Path,
        title: str,
        kind: str,
        *,
        metadata: dict | None = None,
        preview: dict | None = None,
    ) -> Artifact:
        stored_name = path.name
        return self._record(path, stored_name, title, kind, metadata=metadata, preview=preview)

    @staticmethod
    def _register_pdf_font() -> str:
        candidates = [
            Path("C:/Windows/Fonts/arial.ttf"),
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            Path("/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf"),
        ]
        for candidate in candidates:
            if candidate.exists():
                name = "NexaSans"
                if name not in pdfmetrics.getRegisteredFontNames():
                    pdfmetrics.registerFont(TTFont(name, str(candidate)))
                return name
        return "Helvetica"

    @staticmethod
    def _validate_excel(path: Path) -> None:
        workbook = load_workbook(path, read_only=True, data_only=False)
        if "Data" not in workbook.sheetnames or "Metadata" not in workbook.sheetnames:
            raise ValueError("Generated workbook is missing required sheets")
        workbook.close()

    @staticmethod
    def _validate_word(path: Path) -> None:
        document = Document(str(path))
        if not document.paragraphs or not document.core_properties.title:
            raise ValueError("Generated Word document failed structural validation")
        with zipfile.ZipFile(path) as archive:
            if "word/document.xml" not in archive.namelist():
                raise ValueError("Generated Word package is incomplete")

    @staticmethod
    def _validate_powerpoint(path: Path) -> None:
        presentation = Presentation(str(path))
        if len(presentation.slides) < 2:
            raise ValueError("Generated presentation must contain at least two slides")
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.left < 0 or shape.top < 0:
                    raise ValueError("Generated presentation contains off-canvas content")

    @staticmethod
    def _validate_pdf(path: Path) -> None:
        reader = PdfReader(str(path), strict=False)
        if not reader.pages:
            raise ValueError("Generated PDF contains no pages")


def _xml_escape(value: str) -> str:
    return (
        value.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&#39;")
    )


def numeric_summary(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    if not rows:
        return []
    result = []
    for key in rows[0]:
        values = [
            float(row[key])
            for row in rows
            if isinstance(row.get(key), (int, float)) and not isinstance(row.get(key), bool)
        ]
        if values:
            result.append(
                {
                    "field": key,
                    "count": len(values),
                    "minimum": min(values),
                    "maximum": max(values),
                    "average": sum(values) / len(values),
                    "total": sum(values),
                }
            )
    return result
