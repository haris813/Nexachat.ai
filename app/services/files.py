from __future__ import annotations

import csv
import hashlib
import io
import json
import mimetypes
import secrets
import zipfile
from datetime import UTC, datetime, timedelta
from pathlib import Path
from typing import Any

import filetype
from defusedxml.ElementTree import fromstring as safe_xml_fromstring
from docx import Document
from flask import current_app
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
from werkzeug.datastructures import FileStorage
from werkzeug.utils import secure_filename

from ..extensions import db
from ..models import UploadedFile
from .security import SecurityError, safe_display_name

ALLOWED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
    ".txt",
    ".json",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".gif",
    ".mp3",
    ".wav",
    ".m4a",
    ".ogg",
    ".webm",
}

MIME_BY_EXTENSION = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".csv": "text/csv",
    ".txt": "text/plain",
    ".json": "application/json",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".gif": "image/gif",
    ".mp3": "audio/mpeg",
    ".wav": "audio/wav",
    ".m4a": "audio/mp4",
    ".ogg": "audio/ogg",
    ".webm": "audio/webm",
}

OFFICE_MARKERS = {
    ".docx": "word/document.xml",
    ".pptx": "ppt/presentation.xml",
    ".xlsx": "xl/workbook.xml",
}


class FileValidationError(SecurityError):
    pass


def _read_header(path: Path, length: int = 8192) -> bytes:
    with path.open("rb") as handle:
        return handle.read(length)


def _validate_office_zip(path: Path, extension: str) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            marker = OFFICE_MARKERS[extension]
            if marker not in names or "[Content_Types].xml" not in names:
                raise FileValidationError("Office file contents do not match the extension")
            content_types = archive.read("[Content_Types].xml")
            safe_xml_fromstring(content_types)
            total_uncompressed = sum(item.file_size for item in archive.infolist())
            if total_uncompressed > 200 * 1024 * 1024:
                raise FileValidationError("Compressed file expands beyond the safe limit")
    except (zipfile.BadZipFile, KeyError, ValueError) as error:
        raise FileValidationError("The Office file is invalid or corrupted") from error


def validate_file(path: Path, extension: str) -> str:
    extension = extension.lower()
    if extension not in ALLOWED_EXTENSIONS:
        raise FileValidationError(f"Unsupported file type: {extension or 'unknown'}")
    header = _read_header(path)
    if not header:
        raise FileValidationError("Empty files are not supported")
    if extension in OFFICE_MARKERS:
        _validate_office_zip(path, extension)
        return MIME_BY_EXTENSION[extension]
    if extension == ".pdf" and not header.startswith(b"%PDF-"):
        raise FileValidationError("PDF signature does not match the extension")
    if extension in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        try:
            with Image.open(path) as image:
                image.verify()
        except Exception as error:
            raise FileValidationError("Image file is invalid or corrupted") from error
    if extension == ".json":
        try:
            json.loads(path.read_text(encoding="utf-8"))
        except (UnicodeDecodeError, json.JSONDecodeError) as error:
            raise FileValidationError("JSON file is not valid UTF-8 JSON") from error
    if extension in {".txt", ".csv"}:
        try:
            path.read_text(encoding="utf-8")
        except UnicodeDecodeError as error:
            raise FileValidationError("Text and CSV uploads must use UTF-8") from error
    guessed = filetype.guess(path)
    if guessed and extension not in {".txt", ".csv", ".json", ".m4a"}:
        expected_family = MIME_BY_EXTENSION[extension].split("/", 1)[0]
        if guessed.mime.split("/", 1)[0] != expected_family:
            raise FileValidationError("File content does not match the selected extension")
    return MIME_BY_EXTENSION.get(extension, mimetypes.guess_type(path.name)[0] or "application/octet-stream")


class UploadService:
    @staticmethod
    def save(owner_id: str, incoming: FileStorage) -> UploadedFile:
        from .storage import get_upload_storage

        original = safe_display_name(incoming.filename or "", "upload")
        extension = Path(original).suffix.lower()
        safe_stem = secure_filename(Path(original).stem)[:80] or "upload"
        stored_name = f"{secrets.token_hex(16)}-{safe_stem}{extension}"
        owner_dir = (
            Path(current_app.config["UPLOAD_DIR"]) / hashlib.sha256(owner_id.encode()).hexdigest()[:16]
        )
        owner_dir.mkdir(parents=True, exist_ok=True)
        path = owner_dir / stored_name
        incoming.save(path)
        try:
            size = path.stat().st_size
            if size > current_app.config["MAX_CONTENT_LENGTH"]:
                raise FileValidationError("Upload exceeds the configured size limit")
            mime_type = validate_file(path, extension)
            sha256 = hashlib.sha256(path.read_bytes()).hexdigest()
            use_background = bool(
                current_app.config.get("REDIS_URL")
                and size > current_app.config["FILE_BACKGROUND_THRESHOLD_MB"] * 1024 * 1024
            )
            extracted_text, metadata = (
                ("", {"queued": True}) if use_background else extract_file(path, extension)
            )
            storage = get_upload_storage()
            storage_key = storage.save_file(owner_id, stored_name, path)
            expires_at = datetime.now(UTC) + timedelta(hours=current_app.config["FILE_RETENTION_HOURS"])
            record = UploadedFile(
                owner_id=owner_id,
                original_name=original,
                stored_name=stored_name,
                storage_path=storage_key,
                mime_type=mime_type,
                extension=extension,
                size_bytes=size,
                sha256=sha256,
                status="processing" if use_background else "ready",
                extracted_text=extracted_text,
                metadata_json=json.dumps(metadata, ensure_ascii=False),
                expires_at=expires_at,
            )
            db.session.add(record)
            db.session.commit()
            if use_background:
                try:
                    from redis import Redis
                    from rq import Queue

                    queue = Queue("nexachat", connection=Redis.from_url(current_app.config["REDIS_URL"]))
                    queue.enqueue("app.jobs.extract_upload_job", record.id, job_timeout=600)
                except Exception:
                    extracted_text, metadata = extract_file(path, extension)
                    record.status = "ready"
                    record.extracted_text = extracted_text
                    record.metadata_json = json.dumps(metadata, ensure_ascii=False)
                    db.session.commit()
            return record
        except Exception:
            path.unlink(missing_ok=True)
            db.session.rollback()
            raise


def _limit(text: str) -> str:
    return text[: current_app.config["MAX_EXTRACTED_CHARS"]]


def _tabular_preview(rows: list[list[Any]], max_rows: int = 20, max_cols: int = 12) -> dict:
    preview = []
    for row in rows[:max_rows]:
        preview.append([str(value) if value is not None else "" for value in row[:max_cols]])
    return {
        "rows": preview,
        "row_count": len(rows),
        "column_count": max((len(row) for row in rows), default=0),
    }


def extract_file(path: Path, extension: str) -> tuple[str, dict]:
    if extension == ".pdf":
        reader = PdfReader(str(path), strict=False)
        pages = [(page.extract_text() or "").strip() for page in reader.pages]
        return _limit("\n\n".join(pages)), {"pages": len(reader.pages)}
    if extension == ".docx":
        document = Document(str(path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            paragraphs.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        return _limit("\n".join(paragraphs)), {
            "paragraphs": len(document.paragraphs),
            "tables": len(document.tables),
        }
    if extension == ".pptx":
        presentation = Presentation(str(path))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            text = [
                shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
            ]
            slides.append(f"Slide {index}\n" + "\n".join(text))
        return _limit("\n\n".join(slides)), {"slides": len(presentation.slides)}
    if extension == ".xlsx":
        workbook = load_workbook(path, read_only=True, data_only=False)
        parts: list[str] = []
        previews: dict[str, dict] = {}
        for sheet in workbook.worksheets[:10]:
            sheet_rows = [list(row) for row in sheet.iter_rows(values_only=True)]
            previews[sheet.title] = _tabular_preview(sheet_rows)
            parts.append(f"Sheet: {sheet.title}")
            parts.extend(
                "\t".join("" if value is None else str(value) for value in row[:20])
                for row in sheet_rows[:200]
            )
        workbook.close()
        return _limit("\n".join(parts)), {"sheets": workbook.sheetnames, "previews": previews}
    if extension == ".csv":
        csv_text = path.read_text(encoding="utf-8-sig")
        csv_rows = list(csv.reader(io.StringIO(csv_text)))
        return _limit(csv_text), {"preview": _tabular_preview(csv_rows)}
    if extension == ".json":
        payload = json.loads(path.read_text(encoding="utf-8"))
        formatted = json.dumps(payload, ensure_ascii=False, indent=2)
        shape = type(payload).__name__
        count = len(payload) if isinstance(payload, (dict, list)) else 1
        return _limit(formatted), {"root_type": shape, "items": count}
    if extension == ".txt":
        plain_text = path.read_text(encoding="utf-8")
        return _limit(plain_text), {
            "characters": len(plain_text),
            "lines": plain_text.count("\n") + 1,
        }
    if extension in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        with Image.open(path) as image:
            return "", {"width": image.width, "height": image.height, "format": image.format}
    if extension in {".mp3", ".wav", ".m4a", ".ogg", ".webm"}:
        return "", {"audio": True}
    return "", {}
